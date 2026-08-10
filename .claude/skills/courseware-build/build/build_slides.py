#!/usr/bin/env python3
"""Generate the WSQ Budgeting for Small and Medium Enterprises slide deck
(all-white Tertiary house style).

Highly visual: tile grids, flow diagrams, cards and profile cards — no bullet
walls and NO per-step lab slides (the detailed step-by-step lives in the
Learner Guide). Content is driven by course_data.py + data_domain1..7.py so the
deck stays 100% aligned with the LP, LG and labs.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
from data_domain6 import DOMAIN6
from data_domain7 import DOMAIN7
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4 + DOMAIN5 + DOMAIN6 + DOMAIN7

def _find_repo(start):
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")

# ---------------- palette ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,29,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    rect(s,Inches(10.55),Inches(0.72),Inches(2.0),Inches(1.05),BLUE)
    txt(s,Inches(10.55),Inches(0.86),Inches(2.0),Inches(0.5),[[("SME",26,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(10.55),Inches(1.36),Inches(2.0),Inches(0.4),[[("BUDGETING",11,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])
    PAGE["n"]+=1

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        isz=19 if len(ic)<=2 else 13
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,isz,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.9),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.9),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,21,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll produce",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,INK,True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Tools:  ",13,GREY,True),(services,13,GREY,False)],
        [("Detailed step-by-step instructions are in the Learner Guide.",12,GREY,False)]]); footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Check your work",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

def attendance_slide():
    content("Digital Attendance (Mandatory)",[
     "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
     "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
     "Scan the QR code with your mobile phone camera and submit your attendance.",
     "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")

def assessment_slide():
    content("Assessment",[
     C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
     "Format: Open Book — course slides, Learner Guide and approved materials only.",
     C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")

def assessment_flow_slide():
    flow_h("Assessment Flow",[
     "TRAQOM survey — scan the QR code on the LMS",
     "Assessment digital attendance — scan the SSG QR",
     "Sit WA (SAQ) 70 min, then WA (CS) 80 min — open book",
     "Submit your answers on the LMS",
     "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
attendance_slide()
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Qualifications","PhD · ACTA/DACE-certified adult educator"),
  ("Delivers","WSQ courses on budgeting, finance, data analytics and emerging technologies."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with budgeting, accounting or finance (if any).",
 "What you want your business or department budget to achieve after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is stupid.",
 "Mutual respect: agree to disagree.","One conversation at one time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
flow_h("Download Course Material",[
 "Go to lms-tms.tertiaryinfotech.com and log in with your registered email",
 "Open My Courses and select this course",
 "Download the Trainer Slides and Learner Guide (PDF)",
 "Open the Activities / Lab folder for the hands-on templates",
 "Keep the materials handy — the assessment is open book"],kicker="LMS / TMS · lms-tms.tertiaryinfotech.com",color=VIOLET)
two_col("Lesson Plan — 2 Days, 8 hours/day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Topic 1: Introduction to Financial Budgeting (Labs 1–3)",1),
 ("Topic 2: Financial Forecasting (Lab 4)",1),
 ("Topic 3: Budget Preparation (Lab 5)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Topic 4: Budget Control Plan (Lab 6)",1),
 ("Topic 5: Budget Analysis (Labs 7–8)",1),
 ("Topic 6: Budget Approval (Lab 9)",1),
 ("Topic 7: Financial Compliance (Labs 10–12)",1),
 ("Final Assessment: WA (SAQ) + WA (CS)",1),
 ("Daily timing: 9:30am–6:30pm · 1-hour lunch · tea breaks within",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
tile_grid("Skills Framework",[
 ("TSC Title","Budgeting"),
 ("TSC Code","ICT-FIN-3001-1.1"),
 ("Framework","ICT Skills Framework (SkillsFuture Singapore)"),
 ("Coverage","7 abilities (A1–A7) and 10 knowledge areas (K1–K10) mapped across the 7 topics.")],
 kicker="TSC ALIGNMENT",cols=2,size=15)
tile_grid("TSC Abilities (A1–A7)",[a.replace("A"+str(i+1)+": ","") for i,a in enumerate(C.TSC_ABILITIES)],
 kicker="SKILLS FRAMEWORK TSC",cols=2,size=12,icons=[f"A{i}" for i in range(1,8)])
tile_grid("TSC Knowledge (K1–K10)",[k.split(": ",1)[1] for k in C.TSC_KNOWLEDGE],
 kicker="SKILLS FRAMEWORK TSC",cols=2,size=12,icons=[f"K{i}" for i in range(1,11)])
tile_grid("Learning Outcomes",[
 ("LO1 · Analyse","Analyse business strategies and objectives."),
 ("LO2 · Forecast","Carry out financial forecasting."),
 ("LO3 · Prepare","Prepare budget to meet cash flow requirements."),
 ("LO4 · Control","Prepare budget control plan."),
 ("LO5 · Analyse variances","Perform financial analysis to highlight discrepancies."),
 ("LO6 · Report","Report budget and seek approval."),
 ("LO7 · Comply","Perform financial control to ensure compliance.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=14)
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])
assessment_slide()
assessment_flow_slide()
tile_grid("Criteria for Funding",[
 ("75% attendance","Minimum attendance rate of 75% based on the SSG Digital Attendance record."),
 ("Competent","Complete the assessment and be assessed as 'Competent'.")],
 kicker="WSQ FUNDING",cols=1,size=16)

# ---------------- TOPIC EXTRA VISUAL SLIDES ----------------
def topic1_slides():
    content("Activity: Plan Your Vacation",[
     "If you have a budget of $3,000, how would you plan your vacation?",
     "Think: what are your income limits, your big cost items and your trade-offs?",
     "Share your plan — this is exactly what a company budget does at a bigger scale."],kicker="WARM-UP DISCUSSION",size=20)
    tile_grid("What is a Budget?",[
     ("Goals + plan","The goals for the next year(s) plus a detailed plan of how to achieve them."),
     ("Foresight","Identifies potential opportunities and bottlenecks in advance."),
     ("Decision support","Helps management make important decisions when steering the company."),
     ("12–18 months","Companies typically budget when preparing for the next 12–18 months."),
     ("Operational plan","Reduces uncertainty and plans for the future."),
     ("Personal analogy","Saving $1,000 in 12 months needs a plan combining income and expenses — that plan is a budget.")],
     kicker="FOUNDATIONS",cols=2,size=14)
    big_statement("A company's budget is a detailed roadmap.","It shows what needs to be done to obtain the desired results — expressed through the three main financial statements.","COMPANY BUDGET REPRESENTATION")
    cards3("The Three Financial Statements",[
     (BLUE,"Profit & Loss",["Statement of Income","Revenues — the top line","Cost of Goods Sold","SG&A expenses","Shows PROFITABILITY"]),
     (TEAL,"Balance Sheet",["Statement of Financial Position","Assets: cash, receivables, inventory, PP&E","Liabilities: payables, loans","Equity: owners' capital","Shows LIQUIDITY & NET WORTH"]),
     (VIOLET,"Cash Flow Statement",["Cash in and cash out","Operating activities","Investing activities","Financing activities","Shows CASH HEALTH"])],kicker="FINANCIAL STATEMENTS")
    two_col("Is the Budget Standardised?",[
     ("Budgeting — internal",0),("A management-accounting function",1),
     ("Process depends on the company's needs",1),("Customised as required",1),
     ("Accessed by internal users",1)],
     [("Financial statements — external",0),("Governed by IFRS or GAAP",1),
     ("Standard formats for all companies",1),("Used by investors, banks, IRAS",1)],
     kicker="INTERNAL VS EXTERNAL",lhead="Budget",rhead="Financial statements")
    tile_grid("The Goals of a Budget",[
     ("Quantitative plan","Acquiring and using resources over a specified period."),
     ("More surplus","Increase income and reduce expenses to forecast more surplus (disposable income / profit)."),
     ("Track & monitor","Track sales, monitor costs and prepare a cash flow schedule."),
     ("Reduce uncertainty","An operational plan for the future, updated as reality unfolds.")],
     kicker="WHY BUDGET",cols=2,size=14)
    tile_grid("Who Prepares the Budget?",[
     ("Division heads","Own their unit's targets and cost plans."),
     ("Sales teams","Forecast the top line they must deliver."),
     ("Supply chain","Plans purchasing, inventory and logistics costs."),
     ("Marketing","Campaign and promotion spend."),
     ("HR","Headcount and salary planning."),
     ("Finance","Coordinates, consolidates and controls the process.")],
     kicker="EMPLOYEE INVOLVEMENT",cols=2,size=14)
    two_col("Static or Flexible Budget?",[
     ("Static model",0),("Not adjusted throughout the year",1),
     ("Simple, but ignores new circumstances",1)],
     [("Flexible model",0),("Reviewed monthly / quarterly / half-yearly",1),
     ("Forecasts updated with the latest estimates",1),("Lets the company correct course",1)],
     kicker="PLANNING MODELS",lhead="Static",rhead="Flexible (most companies)")
    tile_grid("Key Elements of Budgeting",[
     "A plan expressed in financial terms for attaining an objective.",
     "Prepared and approved before a defined time.",
     "Shows the planned income to be generated.",
     "Shows the probable expenditure to be incurred.",
     "Indicates the capital to be employed during the period."],
     kicker="ESSENTIALS",cols=1,size=14)
    cards3("The Master Budget",[
     (BLUE,"Operating plan",["Revenue budget","Cost of sales","SG&A expenses"]),
     (TEAL,"Investing plan",["Capital expenditure","Fixed assets & depreciation"]),
     (VIOLET,"Financial plan",["Working capital","Financing & interest","Cash plan"])],kicker="AGGREGATING ALL INPUTS")
    big_statement("Master Budget = P&L + Balance Sheet + Cash Flow.","One document aggregating the input of every business unit, department and cost centre for the annual business plan.","CREATING A MASTER BUDGET",color=TEAL)
    tile_grid("Budget Stakeholders",[
     ("Shareholders & Board","More ambitious, less cautious with predictions."),
     ("Mid-management & finance","Prefer forecasts that are easier to achieve."),
     ("Budgeting committee","Reconciles both views and revises the plan."),
     ("Communication","The committee must explain and 'sell' the final numbers to departments.")],
     kicker="WHO SETS THE NUMBERS",cols=2,size=14)
    two_col("Classification of Budget",[
     ("By content",0),("Operating budget — revenue & expenses",1),
     ("Cash budget — inflows & outflows",1),("Capital budget — long-term investments",1),
     ("Financial budget — funding & liquidity",1)],
     [("By method",0),("Baseline — previous plan as the base",1),
     ("Incremental — % or $ on the baseline",1),("Zero-based — starts fresh",1),
     ("Hybrid — a combination",1)],
     kicker="BUDGETING TYPES",lhead="What it covers",rhead="How it is built")

def topic2_slides():
    content("What are Accounting Principles?",[
     "The rules and guidelines companies must follow when reporting financial data.",
     "In the U.S., the FASB issues the standardised set known as GAAP; internationally, IFRS applies.",
     "Budgets are built on the same principles so budget and actuals stay comparable."],kicker="FOUNDATIONS")
    tile_grid("Key Principles of Accounting",[
     ("Accrual","Record revenue when earned, expenses when incurred."),
     ("Conservatism","Recognise expenses and liabilities early, revenue only when assured."),
     ("Consistency","Apply the same methods period after period."),
     ("Full disclosure","Report all information relevant to users."),
     ("Going concern","Assume the business continues to operate."),
     ("Matching","Match expenses to the revenue they generate."),
     ("Materiality","Record items significant to decisions."),
     ("Revenue recognition","Recognise revenue when the performance obligation is met.")],
     kicker="GAAP / IFRS",cols=2,size=13)
    tile_grid("Advantages of Accounting",[
     "Financial information about the business.","Assistance to management.",
     "Replaces memory; facilitates comparison and benchmarking.","Facilitates settlement of tax liabilities.",
     "Facilitates loans and sale of business.","Evidence in court; helps decision making."],
     kicker="WHY IT MATTERS",cols=2,size=14)
    two_col("Accrual vs Cash Method",[
     ("Accrual method",0),("Revenue recorded when EARNED",1),
     ("Expenses recorded when INCURRED",1),("Matches effort to results — required by GAAP/IFRS",1)],
     [("Cash method",0),("Revenue recorded when cash is RECEIVED",1),
     ("Expenses recorded when cash is PAID",1),("Simple — but can mislead on profitability",1)],
     kicker="FINANCIAL ACCOUNTING METHODS",lhead="Accrual",rhead="Cash")
    tile_grid("Key Principles of Corporate Finance",[
     ("What it is","The division dealing with financial and investment decisions."),
     ("Objective","Maximise shareholder value through long- and short-term financial planning."),
     ("Capital investments","Deploy long-term capital through capital budgeting — identify, estimate cash flows, compare, decide."),
     ("Poor budgeting hurts","Excessive or under-funded investment compromises financing costs or operating capacity.")],
     kicker="CORPORATE FINANCE",cols=2,size=14)
    two_col("Capital Financing — Debt vs Equity",[
     ("Debt",0),("Borrow from banks or issue debt securities",1),
     ("Interest is an obligation",1),("Too much debt increases default risk",1)],
     [("Equity",0),("Sell stock to equity investors",1),
     ("Good for long-term expansion funding",1),("Heavy reliance dilutes earnings and value",1)],
     kicker="SOURCING CAPITAL",lhead="Debt",rhead="Equity")
    big_statement("Budgeting is the key to identify and prioritise capital investments.","Capital budgeting compares planned investments with potential proceeds and decides which projects make the capital budget.","CAPITAL INVESTMENTS",color=VIOLET)
    tile_grid("Time Value of Money",[
     ("The core principle","Money available now is worth more than the same amount later, due to its earning capacity."),
     ("In one line","A dollar today is worth more than a dollar tomorrow."),
     ("Why it matters here","Budgets span 1–10 years — long-horizon plans must account for the time value of money."),
     ("How","Discount future cash flows at the required rate of return.")],
     kicker="TVM",cols=2,size=14)
    two_col("Forecasting Horizons & Techniques",[
     ("Horizons",0),("Short term — weeks to months (cash)",1),
     ("Medium term — the budget year",1),("Long term — 3 to 10 years (strategy)",1)],
     [("Techniques",0),("Qualitative — expert judgement, market research",1),
     ("Quantitative — trend, regression, time series",1),
     ("Bottom-up with company parameters · Top-down with macro assumptions",1)],
     kicker="FORECASTING",lhead="How far ahead",rhead="How to forecast")

def topic3_slides():
    flow_h("Financial Budget Process",[
     "Set budget targets & guidelines for the fiscal year",
     "Departments draft their annual plans",
     "Budget review & consolidation",
     "Board review & approved plan",
     "Forecast vs actual — track variance"],kicker="FROM TARGETS TO APPROVED PLAN")
    tile_grid("Broad Budgeting Methods",[
     ("Baseline budget","Begins with a previous plan as the baseline."),
     ("Incremental budget","A % or $ increment on the previous baseline."),
     ("Zero-based budget","Starts fresh like a new plan."),
     ("Hybrid budget","A combination of any of the three above.")],
     kicker="CHOOSE PER ACCOUNT",cols=2,size=15)
    tile_grid("Budgeting Terms",[
     ("Fiscal year","The 12-month period the budget covers."),
     ("Budget targets & guidelines","Top-down goals and rules for the plan."),
     ("Annual plan","The consolidated department submissions."),
     ("Budget / Board review","Review gates before approval."),
     ("Approved plan","The committed budget for the year."),
     ("Forecast · Actual · Variance","Updated estimate · what happened · the difference.")],
     kicker="SPEAK THE LANGUAGE",cols=2,size=13)
    tile_grid("Key Components of a Budget",[
     ("Fixed expenses","Stay the same from month to month — e.g. rent."),
     ("Flexible expenses","Change from month to month — e.g. utilities."),
     ("Total income","Income from operations plus dividends, rental and other sources."),
     ("Disposable income","What remains after taxes — the surplus to allocate.")],
     kicker="BUILDING BLOCKS",cols=2,size=14)
    big_statement("To balance a budget you have 2 choices:","Increase income, or decrease spending.","BUDGETING PRINCIPLE",color=AMBER)
    two_col("Budgeting — Advantages & Success Factors",[
     ("Advantages",0),("Planning orientation & profitability review",1),
     ("Assumptions & performance evaluation",1),("Funding planning & cash allocation",1),
     ("Bottleneck analysis",1)],
     [("Success factors",0),("Management support & employee involvement",1),
     ("Clear organisational goals & responsibility accounting",1),
     ("Flexibility & communication of results",1),("Sound accounting system",1)],
     kicker="MAKE IT WORK",lhead="Why budget",rhead="What makes it succeed")
    tile_grid("The Cash Budget",[
     ("What","An estimation of cash inflows and outflows over a period — weekly, monthly, quarterly or annual."),
     ("Why","Assess whether the entity has sufficient cash to continue operating."),
     ("Insight","Shows cash needs and surpluses for efficient allocation of cash."),
     ("Manage","Sales and expenses must be managed to an optimal level of cash flow.")],
     kicker="CASH IS KING",cols=2,size=14)
    two_col("Short-Term vs Long-Term Cash Budget",[
     ("Short-term",0),("Utility bills, rent, payroll",1),
     ("Payments to suppliers",1),("Operating expenses & investments",1)],
     [("Long-term",0),("Quarterly and annual tax payments",1),
     ("Capital expenditure projects",1),("Long-term investments",1)],
     kicker="TIME FRAMES",lhead="Weeks–months",rhead="Quarters–years")
    two_col("Cash Inflow vs Cash Outflow",[
     ("Cash inflow — money coming in",0),("Sales revenue",1),("Investments",1),("Financing received",1),
     ("Healthy: inflow > outflow",1)],
     [("Cash outflow — money going out",0),("Staff wages & office rent",1),("Supplier payments",1),
     ("Dividends to shareholders",1),("Unhealthy: outflow > inflow",1)],
     kicker="CASH MOVEMENT",lhead="Inflow",rhead="Outflow")
    tile_grid("Working Capital & Cost Behaviour",[
     ("Working capital","Current assets − current liabilities — the measure of liquidity and short-term health."),
     ("Working capital days","Receivable, payable and inventory days drive the cash cycle."),
     ("Fixed cost","Does not vary with volume — rent, salaries, depreciation."),
     ("Variable cost","Varies with volume — materials, commissions, ad spend.")],
     kicker="LIQUIDITY & COSTS",cols=2,size=14)
    flow_h("The 7-Step Master Budget",[
     "Revenue budget","Cost of sales","SG&A budget","Fixed assets","Working capital","Liabilities","Compile"],
     kicker="99 AGENCY CASE ROADMAP",color=TEAL)

def topic4_slides():
    big_statement("Budgetary control: compare, explain, correct.","Prepare budgets, compare them with actual performance, find the reasons for the differences and take corrective actions.","BUDGET CONTROL")
    flow_h("Budgetary Control Process Steps",[
     "Prepare budgets for the future period",
     "Measure actual performance",
     "Compare actuals with the budget standards",
     "Investigate the reasons for differences",
     "Take corrective actions"],kicker="THE CONTROL LOOP")
    cards3("Types of Budget Control",[
     (BLUE,"Operating control",["Covers revenues & operating expenses","Protects day-to-day operations","Targets the desired EBITDA"]),
     (TEAL,"Cash flow control",["Forecast vs actual inflows/outflows","Ensures obligations are always covered","Invests idle cash for returns"]),
     (VIOLET,"Capex control",["Plans large capital expenditures","Only profitable investments proceed","Decisions taken at the right time"])],kicker="THREE CONTROLS")
    tile_grid("Budget Process Road Map & Control (1)",[
     "Develop the budget collaboratively, involving all stakeholders.",
     "Document and communicate policies and procedures; ensure all stakeholders understand their roles.",
     "Understand the users of the budget.",
     "Provide training to those involved in the budget process.",
     "Link cost-management efforts to the budgeting.",
     "Conduct company-wide meetings to establish budget timelines and responsibilities."],
     kicker="GOVERNANCE",cols=2,size=13)
    tile_grid("Budget Process Road Map & Control (2)",[
     "Designate a budget manager to coordinate the process.",
     "Use adequate hardware and software to support budgeting; avoid re-keying data and use templates.",
     "Develop budgets that are flexible for changing conditions.",
     "Ensure the integrity of systems, databases and records.",
     "Make data available electronically to budget users.",
     "Managers are responsible and accountable for their department budgets."],
     kicker="GOVERNANCE",cols=2,size=13)

def topic5_slides():
    tile_grid("Budget · Actual · Variance · Forecast",[
     ("Budget","Sales budgeted at $500,000 for the period."),
     ("Actual","Only $400,000 of sales were made."),
     ("Variance","($100,000) — and ($100,000) ÷ $500,000 = (20%) below budget."),
     ("Report","Dollar and percentage variances appear on the Budget-to-Actual report.")],
     kicker="WORKED EXAMPLE",cols=2,size=14)
    two_col("Budget vs Forecast",[
     ("Budget",0),("A financial plan prepared in advance",1),
     ("The financial expression of the business target",1),("Sets targets · updated annually",1),
     ("Variance analysis: YES",1)],
     [("Forecast",0),("Estimation of future trends from past and present data",1),
     ("Prediction of what the business WILL achieve",1),("No targets · updated at regular intervals",1),
     ("Variance analysis: NO",1)],
     kicker="KNOW THE DIFFERENCE",lhead="Budget",rhead="Forecast")
    two_col("Types of Variances",[
     ("Adverse variance",0),("Actual income LESS than budget",1),
     ("Actual expenditure MORE than budget",1),("A deficit — investigate and correct",1)],
     [("Favourable variance",0),("Actual income MORE than budget",1),
     ("Actual expenditure LESS than budget",1),("A surplus — substantiate and learn",1)],
     kicker="READ THE SIGNS",lhead="Adverse",rhead="Favourable")
    content("Overspend / Underspend Thresholds",[
     "Manage underspends and overspends by setting clear-cut thresholds.",
     "Thresholds can be a dollar value (e.g. ±$5,000) or a percentage (e.g. ±10%).",
     "Variances beyond the threshold trigger investigation and management action.",
     "Both favourable and adverse variances above the threshold must be substantiated."],kicker="VARIANCE MONITORING")
    cards3("Analysis Techniques",[
     (BLUE,"Vertical analysis",["Each line as a % of revenue (P&L)","or of total assets (BS)","Compare structure across periods"]),
     (TEAL,"Horizontal analysis",["Change vs the prior period","In dollars and percentages","Reveals trends over time"]),
     (VIOLET,"Benchmarking",["Compare against industry peers","Identify performance gaps","Set realistic targets"])],kicker="FINANCIAL METHODOLOGIES")
    tile_grid("KPIs for Budget Analysis",[
     ("Income Statement KPIs","Revenue growth, gross margin, EBITDA margin, net margin."),
     ("Balance Sheet KPIs","Working capital, receivable/payable days, asset turnover."),
     ("Liquidity ratios","Current ratio, quick ratio — can we pay short-term obligations?"),
     ("Profitability KPIs","ROA, ROE, return on capital employed."),
     ("Solvency KPIs","Debt-to-equity, interest cover — long-term resilience."),
     ("Financial ratios","Tie every KPI back to the budget assumptions.")],
     kicker="MEASURE WHAT MATTERS",cols=2,size=13)
    tile_grid("Ways to Improve Budgeting",[
     "Keep budgeting and forecasting flexible; implement rolling forecasts.",
     "Budget to your plan and be clear about goals — measurable KPIs and metrics.",
     "Communicate early and often; involve your entire team.",
     "Plan for various scenarios — models and benchmarking.",
     "Track everything; include profit AND cash flow goals.",
     "Let Excel go — adopt purpose-built tools."],
     kicker="BEST PRACTICES",cols=2,size=13)
    tile_grid("Technology & Budgeting",[
     ("ERP modules","SAP, Oracle, NetSuite — holistic budgeting and reporting."),
     ("Budgeting tools","Anaplan, Quicken, Mint — sync directly to the GL or financial system."),
     ("BI dashboards","Tableau, Power BI, SAP BI — valuable data insights."),
     ("Transparency","Common taxonomy and reporting structure — one source of truth.")],
     kicker="BEYOND SPREADSHEETS",cols=2,size=14)
    content("Microsoft Power BI on Xero",[
     "Analyse cash flow, overdue invoices and bills, and profitability by integrating Xero with Power BI.",
     "Power BI imports Xero contacts, invoices, bills and trial balance daily to build a dashboard and reports.",
     "You need a Power BI Pro licence or trial, and the advisor or standard user role in Xero.",
     "Multiple organisations can connect from one Power BI account — each in its own workspace."],kicker="BI DASHBOARDS")

def topic6_slides():
    flow_h("Annual Budget Approval Process",[
     "Departments submit draft budgets",
     "Budget committee reviews & consolidates",
     "Revisions — numbers are 'opened up' and re-sold",
     "Board reviews and approves",
     "Approved plan communicated for tracking"],kicker="FROM DRAFT TO APPROVAL")
    tile_grid("Stakeholder Precautions — the Approval Pack",[
     ("Review of past year","Highlight your successes and changes to last year's plan."),
     ("Goals and objectives","A succinct overview of the basis for your plan."),
     ("Forecast","Rolling monthly/quarterly forecast of plans including costs, in bullets."),
     ("Plan details","Outline key operational areas and specific projects."),
     ("Costs","Realistic costs for each project and operational expense.")],
     kicker="WHAT TO PRESENT",cols=2,size=13)
    tile_grid("Budget Approval — Common Issues",[
     "The approval process takes longer than it should.",
     "The approved plan is not communicated to business units in time to start tracking.",
     "Approval adds no value without a robust tracking tool, framework and resources.",
     "Manual effort makes budgeting tiring — nobody wants to do it.",
     "No transparency on spend without one golden source of truth."],
     kicker="WHAT GOES WRONG",cols=1,size=13)

def topic7_slides():
    tile_grid("Singapore Corporate Taxation",[
     ("Who pays","All companies pay corporate tax under the Income Tax Act on chargeable income derived from Singapore or foreign income remitted into Singapore."),
     ("Flat rate 17%","Levied on chargeable income = taxable revenues − deductible expenses."),
     ("Tax residency","Resident where control and management are exercised — generally where board meetings are held."),
     ("Resident benefits","Start-up exemptions and Avoidance of Double Taxation Agreements (DTAs).")],
     kicker="TAX BASICS",cols=2,size=13)
    tile_grid("Year of Assessment (YA)",[
     ("Definition","The 12-month period in which the company's income is assessed."),
     ("The rule","Income earned in a financial year is assessed in the YA following the financial year end."),
     ("Example","For YA 2018, the basis period is generally 1 Apr 2016 – 31 Mar 2017."),
     ("Why it matters","YA determines exemption eligibility and every filing deadline.")],
     kicker="TIMELINES",cols=2,size=14)
    tile_grid("Start-Up Tax Exemption (first 3 YAs)",[
     ("75% exemption","On the first S$100,000 of chargeable income (from YA 2020)."),
     ("50% exemption","On the next S$100,000 of chargeable income."),
     ("Withholding tax","Payable on certain payments to non-resident companies."),
     ("DTA relief","Tax-residents avoid double taxation on income taxed in DTA partner countries.")],
     kicker="TAX-RESIDENT BENEFITS",cols=2,size=14)
    two_col("ECI & Corporate Tax Filing",[
     ("ECI — Estimated Chargeable Income",0),
     ("An estimate of taxable income for the YA",1),
     ("File within 3 months of the financial year end",1),
     ("Applies to all companies unless exempted",1)],
     [("Corporate tax return",0),
     ("e-File by 30 November (paper phased out from YA 2020)",1),
     ("15 December online concession no longer applies",1),
     ("All corporate tax documents filed annually",1)],
     kicker="FILING DEADLINES",lhead="ECI",rhead="Income tax")
    tile_grid("Penalties for Late or Non-Payment",[
     ("Late filing","An offence — company officers (e.g. directors) may be prosecuted in court."),
     ("Late payment","5% penalty, plus 1% per month the tax remains unpaid (up to 12% in total)."),
     ("Enforcement","IRAS may take further enforcement or legal action to recover unpaid tax."),
     ("Stay compliant","Build tax deadlines into the budget calendar and cash budget.")],
     kicker="CONSEQUENCES",cols=2,size=14)

TOPIC_EXTRAS={1:topic1_slides,2:topic2_slides,3:topic3_slides,4:topic4_slides,5:topic5_slides,6:topic6_slides,7:topic7_slides}

# ---------------- TOPICS + ACTIVITIES ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=f"TSC COVERAGE {t['weighting']}", cols=2, size=13)
    TOPIC_EXTRAS[t["num"]]()
    acts=TOPIC_ACTS[t["num"]]
    if len(acts)>1:
        third=(len(acts)+2)//3
        groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
        while len(groups)<3: groups.append([])
        cards=[]
        for gi,g in enumerate(groups):
            if not g: lbl="—"
            elif len(g)==1: lbl=f"Lab {g[0]['num']}"
            else: lbl=f"Labs {g[0]['num']}–{g[-1]['num']}"
            cards.append((CARD_COLORS[gi], lbl,
                          [a["title"] for a in g] if g else ["—"]))
        cards3(f"Hands-On Activities — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"], kicker=f"TOPIC {t['code']} · HANDS-ON ACTIVITY")
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)
    if t["num"]==3:
        brk("End of Day 1","Day 2 starts at 9:30 am")
    if t["num"]==5:
        brk("Lunch Break","1 hour")

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("Budgeting foundations","Analysed strategies, statements and budget types with Xero (LO1)."),
 ("Financial forecasting","Forecast a full year in Xero Budget Manager from real assumptions (LO2)."),
 ("Budget preparation","Built the 99 Agency master budget through all 7 schedules (LO3)."),
 ("Budget control","Prepared control plans across operating, cash and capex budgets (LO4)."),
 ("Budget analysis","Analysed variances and built a Power BI dashboard on Xero (LO5)."),
 ("Approval & compliance","Reported budgets for approval and applied Singapore tax rules (LO6, LO7).")],
 kicker="LEARNING OUTCOMES",cols=2,size=13)
content("Recommended Courses",[
 "IBF — Financial Analysis for Non-Finance Managers",
 "WSQ — Accounting for Non-Finance Managers",
 "WSQ — Generative AI for Finance and Fintech",
 "WSQ — Unlocking the Power of Accounting and Tax Systems for SMEs",
 "WSQ — Python Programming for Finance"],kicker="KEEP LEARNING")
content("Support",[
 "If you have any enquiries during and after the class, contact us:",
 "Email: enquiry@tertiaryinfotech.com",
 "Tel: +65 6100 0613",
 "Website: www.tertiarycourses.com.sg"],kicker="WE'RE HERE TO HELP")
assessment_slide()
assessment_flow_slide()
attendance_slide()
big_statement("Thank You!","You are now ready to forecast, prepare, control, analyse and report the budget of a small or medium enterprise.","HAPPY BUDGETING",color=TEAL)

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
